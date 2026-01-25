import os
import shutil
from typing import List

from pdf2docx import Converter


def _require_dependency(command: str, friendly_name: str):
    """Ensure an external binary exists before running a conversion."""
    if not shutil.which(command):
        raise RuntimeError(
            f"{friendly_name} is required for this operation but '{command}' was not found on PATH."
        )


class PDFConverter:
    def __init__(self):
        pass

    def pdf_to_word(self, pdf_path: str, output_path: str):
        """
        Convert PDF to Word (DOCX).
        """
        cv = Converter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

    def pdf_to_ppt(self, pdf_path: str, output_path: str):
        """
        Convert PDF to PowerPoint (PPTX) by extracting text and images.
        """
        import io

        import fitz
        from pptx import Presentation
        from pptx.util import Pt

        doc = fitz.open(pdf_path)
        prs = Presentation()

        # Remove default slides if any
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sle[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sle[i]

        for page in doc:
            # Match slide size to page size (approximate)
            page_width = page.rect.width
            page_height = page.rect.height
            prs.slide_width = Pt(page_width)
            prs.slide_height = Pt(page_height)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Extract text blocks
            blocks = page.get_text("dict")["blocks"]

            for b in blocks:
                if b["type"] == 0:  # Text
                    # Bounding box for the whole block
                    bbox = b["bbox"]
                    x, y = bbox[0], bbox[1]
                    w = bbox[2] - bbox[0]
                    h = bbox[3] - bbox[1]

                    txBox = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
                    tf = txBox.text_frame
                    tf.word_wrap = True

                    for line in b["lines"]:
                        p = tf.add_paragraph()
                        for span in line["spans"]:
                            run = p.add_run()
                            run.text = span["text"]
                            run.font.size = Pt(span["size"])
                            # Basic font mapping attempt
                            if "bold" in span["font"].lower():
                                run.font.bold = True
                            if "italic" in span["font"].lower():
                                run.font.italic = True

            # Extract images
            image_list = page.get_images(full=True)
            for img in image_list:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                rects = page.get_image_rects(xref)
                for rect in rects:
                    image_stream = io.BytesIO(image_bytes)
                    try:
                        slide.shapes.add_picture(
                            image_stream,
                            Pt(rect.x0),
                            Pt(rect.y0),
                            width=Pt(rect.width),
                            height=Pt(rect.height),
                        )
                    except Exception:
                        continue

        prs.save(output_path)
        doc.close()

    def pdf_to_excel(self, pdf_path: str, output_path: str):
        """
        Convert PDF to Excel (XLSX) by extracting tables.
        """
        import pandas as pd
        import pdfplumber

        with pdfplumber.open(pdf_path) as pdf:
            all_tables = []
            for i, page in enumerate(pdf.pages):
                # Try multiple strategies to find tables
                found_on_page = []

                # 1. Default strategy
                ts = page.extract_tables()
                if ts:
                    found_on_page.extend(ts)

                # 2. Lines strategy (good for clear grid lines)
                ts_lines = page.extract_tables(
                    {"vertical_strategy": "lines", "horizontal_strategy": "lines"}
                )
                if ts_lines:
                    found_on_page.extend(ts_lines)

                # 3. Text strategy (good for white-space separated tables)
                if not found_on_page:
                    ts_text = page.extract_tables(
                        {"vertical_strategy": "text", "horizontal_strategy": "text"}
                    )
                    if ts_text:
                        found_on_page.extend(ts_text)

                if found_on_page:
                    for table in found_on_page:
                        if not table:
                            continue
                        # Clean table data: remove None, strip whitespace
                        cleaned_table = [
                            [
                                str(cell).strip() if cell is not None else ""
                                for cell in row
                            ]
                            for row in table
                        ]
                        # Remove completely empty rows/columns
                        df = pd.DataFrame(cleaned_table)
                        df.dropna(how="all", axis=0, inplace=True)
                        df.dropna(how="all", axis=1, inplace=True)
                        if not df.empty:
                            all_tables.append(df)
                else:
                    # Fallback: Extract text
                    text = page.extract_text()
                    if text:
                        lines = [
                            line.split() for line in text.split("\n") if line.strip()
                        ]
                        if lines:
                            all_tables.append(pd.DataFrame(lines))

            if not all_tables:
                all_tables.append(pd.DataFrame([["No tables or text found in PDF"]]))

            # Save to Excel
            with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
                for idx, df in enumerate(all_tables):
                    sheet_name = f"Table_{idx+1}"[:31]
                    df.to_excel(
                        writer, sheet_name=sheet_name, index=False, header=False
                    )

    def word_to_pdf(self, word_path: str, output_dir: str) -> str:
        """
        Convert Word (DOC/DOCX) to PDF using LibreOffice.
        Returns the path to the created PDF.
        """
        import subprocess

        _require_dependency("libreoffice", "LibreOffice")

        # LibreOffice converts to the same directory
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            word_path,
        ]

        try:
            subprocess.run(
                cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )

            # Construct expected output path
            filename = os.path.basename(word_path)
            name_without_ext = os.path.splitext(filename)[0]
            pdf_filename = f"{name_without_ext}.pdf"
            return os.path.join(output_dir, pdf_filename)

        except subprocess.CalledProcessError as e:
            raise Exception(f"LibreOffice conversion failed: {e.stderr.decode()}")

    def ppt_to_pdf(self, ppt_path: str, output_dir: str) -> str:
        """
        Convert PowerPoint (PPT/PPTX) to PDF using LibreOffice.
        Returns the path to the created PDF.
        """
        import subprocess

        _require_dependency("libreoffice", "LibreOffice")

        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            ppt_path,
        ]

        try:
            subprocess.run(
                cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )

            filename = os.path.basename(ppt_path)
            name_without_ext = os.path.splitext(filename)[0]
            pdf_filename = f"{name_without_ext}.pdf"
            return os.path.join(output_dir, pdf_filename)

        except subprocess.CalledProcessError as e:
            raise Exception(f"LibreOffice conversion failed: {e.stderr.decode()}")

    def excel_to_pdf(self, excel_path: str, output_dir: str) -> str:
        """
        Convert Excel (XLS/XLSX) to PDF using LibreOffice.
        Returns the path to the created PDF.
        """
        import subprocess

        _require_dependency("libreoffice", "LibreOffice")

        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            excel_path,
        ]

        try:
            subprocess.run(
                cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )

            filename = os.path.basename(excel_path)
            name_without_ext = os.path.splitext(filename)[0]
            pdf_filename = f"{name_without_ext}.pdf"
            return os.path.join(output_dir, pdf_filename)

        except subprocess.CalledProcessError as e:
            raise Exception(f"LibreOffice conversion failed: {e.stderr.decode()}")

    def pdf_to_jpg(self, pdf_path: str, output_dir: str) -> List[str]:
        """
        Convert PDF pages to JPG images.
        Returns list of paths to created images.
        """
        import fitz

        doc = fitz.open(pdf_path)
        output_files = []

        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
            output_filename = f"page_{i+1}_{os.path.basename(pdf_path)}.jpg"
            output_path = os.path.join(output_dir, output_filename)
            pix.save(output_path)
            output_files.append(output_path)

        doc.close()
        return output_files

    def jpg_to_pdf(self, image_paths: List[str], output_path: str):
        """
        Convert JPG/PNG images to a single PDF.
        """
        import fitz

        doc = fitz.open()

        for img_path in image_paths:
            img = fitz.open(img_path)
            pdfbytes = img.convert_to_pdf()
            img.close()

            imgPDF = fitz.open("pdf", pdfbytes)
            page_src = imgPDF[0]
            rect = page_src.rect

            page = doc.new_page(width=rect.width, height=rect.height)
            page.show_pdf_page(rect, imgPDF, 0)

        doc.save(output_path)
        doc.close()

    def html_to_pdf(self, html_path: str, output_path: str):
        """
        Convert HTML to PDF.
        """
        import fitz

        doc = fitz.open(html_path)
        pdfbytes = doc.convert_to_pdf()
        with open(output_path, "wb") as f:
            f.write(pdfbytes)
        doc.close()

    def pdf_to_pdfa(self, pdf_path: str, output_path: str):
        """
        Convert PDF to PDF/A using Ghostscript.
        """
        import subprocess

        _require_dependency("gs", "Ghostscript")

        # Try PDF/A-2b
        cmd = [
            "gs",
            "-dPDFA=2",
            "-dBATCH",
            "-dNOPAUSE",
            "-sColorConversionStrategy=RGB",
            "-sProcessColorModel=DeviceRGB",
            "-sDEVICE=pdfwrite",
            "-sPDFACompatibilityPolicy=1",
            f"-sOutputFile={output_path}",
            pdf_path,
        ]

        try:
            subprocess.run(
                cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )
        except subprocess.CalledProcessError as e:
            print(f"Ghostscript failed: {e.stderr.decode()}")
            # Fallback to simple pdfwrite if PDF/A fails (e.g. missing color profile)
            print("Falling back to standard PDF conversion...")
            cmd_fallback = [
                "gs",
                "-dBATCH",
                "-dNOPAUSE",
                "-sDEVICE=pdfwrite",
                f"-sOutputFile={output_path}",
                pdf_path,
            ]
            subprocess.run(
                cmd_fallback, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )

    def scan_to_pdf(
        self, image_paths: List[str], output_path: str, enhance: bool = True
    ):
        """
        Convert scanned images to PDF, optionally enhancing them.
        """
        import fitz
        from PIL import Image, ImageEnhance

        doc = fitz.open()

        for img_path in image_paths:
            if enhance:
                # Enhance image using Pillow
                try:
                    img: Image.Image = Image.open(img_path)
                    # Convert to grayscale
                    img = img.convert("L")
                    # Increase contrast
                    enhancer = ImageEnhance.Contrast(img)
                    img = enhancer.enhance(1.5)

                    # Save to temp file
                    temp_img_path = f"{img_path}_enhanced.jpg"
                    img.save(temp_img_path)

                    img_doc = fitz.open(temp_img_path)
                    pdfbytes = img_doc.convert_to_pdf()
                    img_pdf = fitz.open("pdf", pdfbytes)
                    doc.insert_pdf(img_pdf)
                    img_doc.close()
                    img_pdf.close()

                    # Cleanup temp file
                    if os.path.exists(temp_img_path):
                        os.remove(temp_img_path)

                except Exception as e:
                    print(f"Failed to enhance image {img_path}: {e}")
                    # Fallback to original
                    img_doc = fitz.open(img_path)
                    pdfbytes = img_doc.convert_to_pdf()
                    img_pdf = fitz.open("pdf", pdfbytes)
                    doc.insert_pdf(img_pdf)
                    img_doc.close()
                    img_pdf.close()
            else:
                img_doc = fitz.open(img_path)
                pdfbytes = img_doc.convert_to_pdf()
                img_pdf = fitz.open("pdf", pdfbytes)
                doc.insert_pdf(img_pdf)
                img_doc.close()
                img_pdf.close()

        doc.save(output_path)
        doc.close()

    def ocr_pdf(self, pdf_path: str, output_path: str, language: str = "eng"):
        """
        OCR PDF using pytesseract (Tesseract).
        """
        import io

        import fitz
        import pytesseract
        from PIL import Image

        _require_dependency("tesseract", "Tesseract OCR")

        doc = fitz.open(pdf_path)
        out_doc = fitz.open()

        for page in doc:
            # Get image from page
            pix = page.get_pixmap()
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))

            # OCR to PDF
            try:
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                    img, extension="pdf", lang=language
                )
                img_pdf = fitz.open("pdf", pdf_bytes)
                out_doc.insert_pdf(img_pdf)
                img_pdf.close()
            except Exception as e:
                print(f"OCR failed for page: {e}")
                # Fallback: just insert original page (as image or original)
                # If we insert original page, it might not be searchable if it was image-only.
                # But better than failing.
                out_doc.insert_pdf(doc, from_page=page.number, to_page=page.number)

        out_doc.save(output_path)
        out_doc.close()
        doc.close()

    def pdf_to_markdown(self, pdf_path: str, output_path: str):
        """
        Convert PDF to Markdown format.
        Extracts text with formatting, headings, lists, and tables.
        """
        import fitz

        doc = fitz.open(pdf_path)
        markdown_lines = []

        for page_num, page in enumerate(doc, start=1):
            # Get text as blocks for structure detection
            blocks = page.get_text("dict")["blocks"]

            for block in blocks:
                if block["type"] == 0:  # Text block
                    # Process each line and its spans
                    if "lines" in block:
                        for line in block["lines"]:
                            if "spans" in line and line["spans"]:
                                max_font = max(span.get("size", 12) for span in line["spans"])
                                line_text = "".join(span["text"] for span in line["spans"])

                                if not line_text.strip():
                                    continue

                                # Detect headings based on font size
                                if max_font > 18:
                                    markdown_lines.append(f"# {line_text}\n")
                                elif max_font > 16:
                                    markdown_lines.append(f"## {line_text}\n")
                                elif max_font > 14:
                                    markdown_lines.append(f"### {line_text}\n")
                                else:
                                    # Check for list items
                                    if line_text.strip().startswith(("-", "•", "*")):
                                        markdown_lines.append(f"{line_text}\n")
                                    else:
                                        markdown_lines.append(f"{line_text}\n")

            # Page separator
            markdown_lines.append("\n---\n\n")

        doc.close()

        # Write to file
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("".join(markdown_lines))

    def pdf_to_txt(self, pdf_path: str, output_path: str):
        """
        Convert PDF to plain text format.
        Simple text extraction without formatting.
        """
        import fitz

        doc = fitz.open(pdf_path)
        full_text = []

        for page in doc:
            text = page.get_text("text")
            full_text.append(text)

        doc.close()

        # Write to file
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n\n".join(full_text))

    def pdf_to_epub(self, pdf_path: str, output_path: str):
        """
        Convert PDF to EPUB format.
        Creates chapters from PDF pages.
        """
        import fitz
        from ebooklib import epub

        doc = fitz.open(pdf_path)
        book = epub.EpubBook()

        # Set metadata
        book.set_identifier("pdfsmarteditor-" + str(hash(pdf_path)))
        book.set_title(os.path.splitext(os.path.basename(pdf_path))[0])
        book.set_language("en")

        # Create chapters from pages
        chapters = []
        toc = []

        for page_num, page in enumerate(doc, start=1):
            # Extract text and images from page
            text = page.get_text("text")
            html_content = f"<h1>Page {page_num}</h1>\n"

            # Add text as paragraphs
            for para in text.split("\n\n"):
                if para.strip():
                    html_content += f"<p>{para}</p>\n"

            # Extract images
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                if base_image:
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    img_filename = f"page{page_num}_img{img_index}.{image_ext}"

                    # Add image to EPUB
                    img_item = epub.EpubItem(
                        uid=f"img_{page_num}_{img_index}",
                        file_name=f"images/{img_filename}",
                        media_type=f"image/{image_ext}",
                        content=image_bytes,
                    )
                    book.add_item(img_item)
                    html_content += f'<img src="images/{img_filename}"/>\n'

            # Create chapter
            chapter = epub.EpubHtml(
                title=f"Page {page_num}",
                file_name=f"chap_{page_num}.xhtml",
                content=html_content,
            )
            book.add_item(chapter)
            chapters.append(chapter)
            toc.append(chapter)

        # Add navigation
        book.toc = tuple(toc)
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())

        # Add CSS
        style = "body { font-family: Times, serif; }"
        nav_css = epub.EpubItem(
            uid="style_nav",
            file_name="style/nav.css",
            media_type="text/css",
            content=style,
        )
        book.add_item(nav_css)

        # Create spine
        book.spine = ["nav"] + chapters

        # Write EPUB
        epub.write_epub(output_path, book, {})
        doc.close()

    def pdf_to_svg(self, pdf_path: str, output_dir: str) -> List[str]:
        """
        Convert PDF pages to SVG format.
        Returns list of paths to created SVG files.
        """
        import fitz

        doc = fitz.open(pdf_path)
        output_files = []

        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)

        for page_num, page in enumerate(doc, start=1):
            # Get SVG content using get_svg_image
            svg = page.get_svg_image()
            svg_data = svg.tobytes("utf-8") if isinstance(svg, bytes) else svg
            output_filename = f"page_{page_num}_{os.path.basename(pdf_path)}.svg"
            output_path = os.path.join(output_dir, output_filename)

            # Write SVG
            with open(output_path, "w", encoding="utf-8") as f:
                if isinstance(svg_data, bytes):
                    f.write(svg_data.decode("utf-8"))
                else:
                    f.write(svg_data)

            output_files.append(output_path)

        doc.close()
        return output_files

    def markdown_to_pdf(self, md_path: str, output_path: str):
        """
        Convert Markdown to PDF.
        Uses reportlab for PDF generation.
        """
        import markdown
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak

        # Read markdown file
        with open(md_path, "r", encoding="utf-8") as f:
            md_content = f.read()

        # Convert markdown to HTML
        html = markdown.markdown(
            md_content, extensions=["tables", "fenced_code", "codehilite"]
        )

        # Parse HTML and create PDF
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Create custom styles (unique names to avoid conflicts)
        h1_style = ParagraphStyle(
            name="CustomH1",
            parent=styles["Heading1"],
            fontSize=18,
            textColor="#000000",
        )
        h2_style = ParagraphStyle(
            name="CustomH2",
            parent=styles["Heading2"],
            fontSize=14,
            textColor="#000000",
        )
        h3_style = ParagraphStyle(
            name="CustomH3",
            parent=styles["Heading3"],
            fontSize=12,
            textColor="#000000",
        )
        code_style = ParagraphStyle(
            name="CustomCode",
            parent=styles["Code"],
            fontName="Courier",
            fontSize=9,
            leftIndent=20,
            backColor="#f5f5f5",
        )

        # Simple HTML to PDF conversion
        import re

        # Split by HTML tags for basic parsing
        parts = re.split(r"<(h[1-6]|p|pre|code|li|ul|ol|/[^>]+)>", html)

        current_tag = "p"
        for part in parts:
            part = part.strip()
            if not part:
                continue

            if part in ("h1", "h2", "h3", "h4", "h5", "h6", "p", "pre", "code", "li"):
                current_tag = part
            elif part.startswith("/"):
                current_tag = "p"
            else:
                # Decode HTML entities
                text = part.replace("&lt;", "<").replace("&gt;", ">").replace("&amp;", "&")

                if current_tag == "h1":
                    story.append(Paragraph(text, h1_style))
                    story.append(Spacer(1, 0.2 * inch))
                elif current_tag == "h2":
                    story.append(Paragraph(text, h2_style))
                    story.append(Spacer(1, 0.15 * inch))
                elif current_tag == "h3":
                    story.append(Paragraph(text, h3_style))
                    story.append(Spacer(1, 0.1 * inch))
                elif current_tag in ("pre", "code"):
                    story.append(Paragraph(text, code_style))
                    story.append(Spacer(1, 0.1 * inch))
                else:
                    story.append(Paragraph(text, styles["BodyText"]))
                    story.append(Spacer(1, 0.1 * inch))

        doc.build(story)

    def txt_to_pdf(self, txt_path: str, output_path: str):
        """
        Convert plain text to PDF.
        Creates a simple PDF with the text content.
        """
        import fitz

        # Read text file
        with open(txt_path, "r", encoding="utf-8") as f:
            text_content = f.read()

        # Create PDF
        doc = fitz.open()
        page = doc.new_page(width=595, height=842)  # A4 size

        # Insert text
        # Fit text to page with margins
        margin = 50
        rect = fitz.Rect(margin, margin, 595 - margin, 842 - margin)

        # Split text into lines that fit the page width
        import textwrap

        lines = text_content.split("\n")
        wrapped_lines = []
        for line in lines:
            wrapped_lines.extend(textwrap.wrap(line, width=80))

        # Create text pages
        y_position = margin
        line_height = 14

        for line in wrapped_lines:
            if y_position > 842 - margin:
                # New page
                page = doc.new_page(width=595, height=842)
                y_position = margin

            page.insert_text(
                fitz.Point(margin, y_position), line, fontsize=11, fontname="helvetica"
            )
            y_position += line_height

        doc.save(output_path)
        doc.close()

    def csv_to_pdf(self, csv_path: str, output_path: str):
        """
        Convert CSV to PDF with formatted table.
        """
        import fitz
        import pandas as pd

        # Read CSV
        df = pd.read_csv(csv_path)

        # Create PDF
        doc = fitz.open()
        page = doc.new_page(width=595, height=842)  # A4

        # Table settings
        margin = 50
        col_width = 100
        row_height = 25
        y_position = margin

        # Draw headers
        for col_idx, col_name in enumerate(df.columns):
            x_position = margin + col_idx * col_width

            # Draw header background
            rect = fitz.Rect(x_position, y_position, x_position + col_width, y_position + row_height)
            page.draw_rect(rect, color=(0.7, 0.7, 0.7), fill=(0.7, 0.7, 0.7))

            # Draw header text
            page.insert_text(
                fitz.Point(x_position + 5, y_position + 17),
                str(col_name),
                fontsize=10,
                fontname="helvetica-bold",
            )

        y_position += row_height

        # Draw data rows
        for _, row in df.iterrows():
            if y_position > 842 - margin - row_height:
                # New page
                page = doc.new_page(width=595, height=842)
                y_position = margin

            # Redraw headers on new page
            for col_idx, col_name in enumerate(df.columns):
                x_position = margin + col_idx * col_width
                rect = fitz.Rect(x_position, y_position, x_position + col_width, y_position + row_height)
                page.draw_rect(rect, color=(0.7, 0.7, 0.7), fill=(0.7, 0.7, 0.7))
                page.insert_text(
                    fitz.Point(x_position + 5, y_position + 17),
                    str(col_name),
                    fontsize=10,
                    fontname="helvetica-bold",
                )
            y_position += row_height

            # Draw row data
            for col_idx, value in enumerate(row):
                x_position = margin + col_idx * col_width

                # Draw cell border
                rect = fitz.Rect(x_position, y_position, x_position + col_width, y_position + row_height)
                page.draw_rect(rect, color=(0.5, 0.5, 0.5))

                # Draw cell text (truncate if too long)
                text = str(value)[:20] + "..." if len(str(value)) > 20 else str(value)
                page.insert_text(
                    fitz.Point(x_position + 5, y_position + 17),
                    text,
                    fontsize=9,
                    fontname="helvetica",
                )

            y_position += row_height

        doc.save(output_path)
        doc.close()

    def json_to_pdf(self, json_path: str, output_path: str):
        """
        Convert JSON to PDF.
        Handles arrays of objects as formatted tables.
        """
        import fitz
        import json

        # Read JSON
        with open(json_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        # Create PDF
        doc = fitz.open()
        page = doc.new_page(width=595, height=842)  # A4

        # Table settings
        margin = 50
        col_width = 100
        row_height = 25
        y_position = margin

        # Handle array of objects
        if isinstance(data, list):
            if not data:
                # Empty array
                page.insert_text(
                    fitz.Point(margin, y_position),
                    "Empty JSON array",
                    fontsize=12,
                    fontname="helvetica",
                )
            else:
                # Get all unique keys from all objects
                all_keys = set()
                for item in data:
                    if isinstance(item, dict):
                        all_keys.update(item.keys())
                columns = list(all_keys)

                # Draw headers
                for col_idx, col_name in enumerate(columns):
                    x_position = margin + col_idx * col_width

                    # Draw header background
                    rect = fitz.Rect(x_position, y_position, x_position + col_width, y_position + row_height)
                    page.draw_rect(rect, color=(0.7, 0.7, 0.7), fill=(0.7, 0.7, 0.7))

                    # Draw header text
                    page.insert_text(
                        fitz.Point(x_position + 5, y_position + 17),
                        str(col_name),
                        fontsize=10,
                        fontname="helvetica-bold",
                    )

                y_position += row_height

                # Draw data rows
                for row_data in data:
                    if y_position > 842 - margin - row_height:
                        # New page
                        page = doc.new_page(width=595, height=842)
                        y_position = margin

                        # Redraw headers
                        for col_idx, col_name in enumerate(columns):
                            x_position = margin + col_idx * col_width
                            rect = fitz.Rect(x_position, y_position, x_position + col_width, y_position + row_height)
                            page.draw_rect(rect, color=(0.7, 0.7, 0.7), fill=(0.7, 0.7, 0.7))
                            page.insert_text(
                                fitz.Point(x_position + 5, y_position + 17),
                                str(col_name),
                                fontsize=10,
                                fontname="helvetica-bold",
                            )
                        y_position += row_height

                    # Draw row data
                    for col_idx, key in enumerate(columns):
                        x_position = margin + col_idx * col_width

                        # Draw cell border
                        rect = fitz.Rect(x_position, y_position, x_position + col_width, y_position + row_height)
                        page.draw_rect(rect, color=(0.5, 0.5, 0.5))

                        # Get value
                        value = row_data.get(key, "") if isinstance(row_data, dict) else ""
                        text = str(value)[:20] + "..." if len(str(value)) > 20 else str(value)

                        # Draw cell text
                        page.insert_text(
                            fitz.Point(x_position + 5, y_position + 17),
                            text,
                            fontsize=9,
                            fontname="helvetica",
                        )

                    y_position += row_height

        elif isinstance(data, dict):
            # Handle single object - pretty print
            import textwrap

            for key, value in data.items():
                if y_position > 842 - margin - row_height:
                    page = doc.new_page(width=595, height=842)
                    y_position = margin

                # Draw key-value pair
                text = f"{key}: {value}"
                wrapped_lines = textwrap.wrap(text, width=70)

                for line in wrapped_lines:
                    page.insert_text(
                        fitz.Point(margin, y_position),
                        line,
                        fontsize=11,
                        fontname="helvetica",
                    )
                    y_position += row_height

                y_position += 5
        else:
            # Handle primitive value
            page.insert_text(
                fitz.Point(margin, y_position),
                str(data),
                fontsize=12,
                fontname="helvetica",
            )

        doc.save(output_path)
        doc.close()

    def batch_convert(
        self, files: List[str], output_dir: str, conversion_type: str
    ) -> List[str]:
        """
        Batch convert multiple files using the same conversion type.

        Args:
            files: List of input file paths
            output_dir: Directory for output files
            conversion_type: Type of conversion (e.g., 'pdf-to-word', 'word-to-pdf')

        Returns:
            List of output file paths
        """
        import uuid

        output_files = []

        # Map conversion types to methods
        conversion_map = {
            "pdf-to-word": self.pdf_to_word,
            "pdf-to-ppt": self.pdf_to_ppt,
            "pdf-to-excel": self.pdf_to_excel,
            "pdf-to-jpg": self.pdf_to_jpg,
            "pdf-to-markdown": self.pdf_to_markdown,
            "pdf-to-txt": self.pdf_to_txt,
            "pdf-to-epub": self.pdf_to_epub,
            "pdf-to-svg": self.pdf_to_svg,
            "word-to-pdf": self.word_to_pdf,
            "ppt-to-pdf": self.ppt_to_pdf,
            "excel-to-pdf": self.excel_to_pdf,
            "html-to-pdf": self.html_to_pdf,
            "markdown-to-pdf": self.markdown_to_pdf,
            "txt-to-pdf": self.txt_to_pdf,
            "csv-to-pdf": self.csv_to_pdf,
            "json-to-pdf": self.json_to_pdf,
        }

        converter_func = conversion_map.get(conversion_type)

        if not converter_func:
            raise ValueError(f"Unsupported conversion type: {conversion_type}")

        for file_path in files:
            try:
                # Determine output extension
                ext_map = {
                    "pdf-to-word": ".docx",
                    "pdf-to-ppt": ".pptx",
                    "pdf-to-excel": ".xlsx",
                    "pdf-to-markdown": ".md",
                    "pdf-to-txt": ".txt",
                    "pdf-to-epub": ".epub",
                    "word-to-pdf": ".pdf",
                    "ppt-to-pdf": ".pdf",
                    "excel-to-pdf": ".pdf",
                    "html-to-pdf": ".pdf",
                    "markdown-to-pdf": ".pdf",
                    "txt-to-pdf": ".pdf",
                    "csv-to-pdf": ".pdf",
                    "json-to-pdf": ".pdf",
                }

                ext = ext_map.get(conversion_type, ".pdf")
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_path = os.path.join(output_dir, f"{base_name}_converted{ext}")

                # Call the conversion function
                if conversion_type == "pdf-to-jpg":
                    # Returns list of files
                    result = converter_func(file_path, output_dir)
                    output_files.extend(result)
                elif conversion_type == "pdf-to-svg":
                    # Returns list of files
                    result = converter_func(file_path, output_dir)
                    output_files.extend(result)
                elif conversion_type in ["word-to-pdf", "ppt-to-pdf", "excel-to-pdf"]:
                    # Returns output path
                    result = converter_func(file_path, output_dir)
                    output_files.append(result)
                else:
                    # Standard conversion (input_path, output_path)
                    converter_func(file_path, output_path)
                    output_files.append(output_path)

            except Exception as e:
                print(f"Failed to convert {file_path}: {e}")
                # Continue with other files
                continue

        return output_files

    def merge_folder(self, folder_path: str, output_path: str) -> int:
        """
        Merge all PDFs from a folder into a single PDF.

        Args:
            folder_path: Path to folder containing PDFs
            output_path: Path for merged output PDF

        Returns:
            Number of PDFs merged
        """
        import glob

        # Find all PDFs in folder
        pdf_pattern = os.path.join(folder_path, "*.pdf")
        pdf_files = sorted(glob.glob(pdf_pattern))

        if not pdf_files:
            raise ValueError(f"No PDF files found in {folder_path}")

        if len(pdf_files) < 2:
            raise ValueError("At least 2 PDF files are required for merging")

        # Merge PDFs using existing merge functionality
        from pdfsmarteditor.core.manipulator import PDFManipulator

        manipulator = PDFManipulator()
        manipulator.merge_pdfs(pdf_files, output_path)

        return len(pdf_files)

    def apply_template(
        self, template: dict, files: List[str], output_dir: str
    ) -> List[str]:
        """
        Apply template settings to multiple files.

        Template options:
        - watermark: Add watermark text
        - rotate: Rotate pages by degrees
        - compress: Compression level (0-9)
        - protect: Add password protection

        Args:
            template: Dictionary with template settings
            files: List of input PDF file paths
            output_dir: Directory for output files

        Returns:
            List of output file paths
        """
        import uuid

        output_files = []
        from pdfsmarteditor.core.manipulator import PDFManipulator

        manipulator = PDFManipulator()

        for file_path in files:
            try:
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_path = os.path.join(output_dir, f"{base_name}_processed.pdf")

                # Apply transformations in order
                current_path = file_path

                # Watermark
                if "watermark" in template:
                    temp_path = os.path.join(output_dir, f"temp_{uuid.uuid4()}.pdf")
                    wm = template["watermark"]
                    manipulator.add_watermark(
                        current_path,
                        wm.get("text", ""),
                        temp_path,
                        wm.get("opacity", 0.3),
                        wm.get("rotation", 45),
                        wm.get("font_size", 50),
                        wm.get("color", (0, 0, 0)),
                    )
                    if current_path != file_path:
                        os.remove(current_path)
                    current_path = temp_path

                # Rotate
                if "rotate" in template:
                    temp_path = os.path.join(output_dir, f"temp_{uuid.uuid4()}.pdf")
                    manipulator.rotate_pdf(
                        current_path, temp_path, template["rotate"], None
                    )
                    if current_path != file_path:
                        os.remove(current_path)
                    current_path = temp_path

                # Compress
                if "compress" in template:
                    temp_path = os.path.join(output_dir, f"temp_{uuid.uuid4()}.pdf")
                    manipulator.compress_pdf(
                        current_path, temp_path, template["compress"]
                    )
                    if current_path != file_path:
                        os.remove(current_path)
                    current_path = temp_path

                # Protect
                if "protect" in template:
                    temp_path = os.path.join(output_dir, f"temp_{uuid.uuid4()}.pdf")
                    manipulator.protect_pdf(current_path, template["protect"], temp_path)
                    if current_path != file_path:
                        os.remove(current_path)
                    current_path = temp_path

                # Move to final output
                if current_path != output_path:
                    import shutil

                    shutil.move(current_path, output_path)

                output_files.append(output_path)

            except Exception as e:
                print(f"Failed to apply template to {file_path}: {e}")
                continue

        return output_files
