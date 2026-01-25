import os
import uuid
from typing import Iterator

import docx
import fitz
import pandas as pd
import pytest
from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation

from api.deps import TEMP_DIR
from api.main import app


def _build_pdf(path: str, pages: int = 1) -> str:
    doc = fitz.open()
    for i in range(pages):
        page = doc.new_page()
        page.insert_text((72, 72), f"Page {i + 1}")
    doc.save(path)
    doc.close()
    return path


def _build_docx(path: str) -> str:
    document = docx.Document()
    document.add_paragraph("This is a DOCX test file.")
    document.save(path)
    return path


def _build_pptx(path: str) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test PPTX"
    prs.save(path)
    return path


def _build_excel(path: str) -> str:
    df = pd.DataFrame({"A": [1, 2], "B": [3, 4]})
    df.to_excel(path, index=False)
    return path


def _build_image(path: str) -> str:
    img = Image.new("RGB", (100, 100), color=(255, 0, 0))
    img.save(path)
    return path


@pytest.fixture(scope="session")
def api_client() -> Iterator[TestClient]:
    yield TestClient(app)


@pytest.fixture
def sample_pdf(tmp_path) -> Iterator[str]:
    file_path = tmp_path / f"sample_{uuid.uuid4().hex}.pdf"
    _build_pdf(str(file_path), pages=1)
    yield str(file_path)


@pytest.fixture
def multi_page_pdf(tmp_path) -> Iterator[str]:
    file_path = tmp_path / f"multi_{uuid.uuid4().hex}.pdf"
    _build_pdf(str(file_path), pages=3)
    yield str(file_path)


@pytest.fixture
def sample_docx(tmp_path) -> Iterator[str]:
    file_path = tmp_path / f"doc_{uuid.uuid4().hex}.docx"
    _build_docx(str(file_path))
    yield str(file_path)


@pytest.fixture
def sample_pptx(tmp_path) -> Iterator[str]:
    file_path = tmp_path / f"ppt_{uuid.uuid4().hex}.pptx"
    _build_pptx(str(file_path))
    yield str(file_path)


@pytest.fixture
def sample_excel(tmp_path) -> Iterator[str]:
    file_path = tmp_path / f"excel_{uuid.uuid4().hex}.xlsx"
    _build_excel(str(file_path))
    yield str(file_path)


@pytest.fixture
def sample_image(tmp_path) -> Iterator[str]:
    file_path = tmp_path / f"img_{uuid.uuid4().hex}.png"
    _build_image(str(file_path))
    yield str(file_path)


@pytest.fixture
def sample_html(tmp_path) -> Iterator[str]:
    file_path = tmp_path / f"page_{uuid.uuid4().hex}.html"
    with open(file_path, "w") as fh:
        fh.write("<html><body><p>HTML to PDF</p></body></html>")
    yield str(file_path)


# =============================================================================
# Phase 1-3 Additional Fixtures
# =============================================================================

@pytest.fixture
def blank_pdf(tmp_path) -> Iterator[str]:
    """Create a blank PDF for testing blank page removal."""
    file_path = tmp_path / f"blank_{uuid.uuid4().hex}.pdf"
    doc = fitz.open()
    doc.save(str(file_path))
    doc.close()
    yield str(file_path)


@pytest.fixture
def pdf_with_annotations(tmp_path) -> Iterator[str]:
    """Create a PDF with annotations for testing."""
    file_path = tmp_path / f"annotated_{uuid.uuid4().hex}.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Test Content")
    # Add a highlight annotation
    annot = page.add_rect_annot(
        fitz.Rect(100, 100, 200, 150),
        "Highlight",
        author="Test User",
    )
    doc.save(str(file_path))
    doc.close()
    yield str(file_path)


@pytest.fixture
def sample_markdown(tmp_path) -> Iterator[str]:
    """Create a Markdown file for testing."""
    file_path = tmp_path / f"doc_{uuid.uuid4().hex}.md"
    with open(file_path, "w") as fh:
        fh.write("# Heading\n\nThis is a **test** Markdown file.\n\n- Item 1\n- Item 2\n")
    yield str(file_path)


@pytest.fixture
def sample_txt(tmp_path) -> Iterator[str]:
    """Create a plain text file for testing."""
    file_path = tmp_path / f"doc_{uuid.uuid4().hex}.txt"
    with open(file_path, "w") as fh:
        fh.write("This is a plain text file for testing TXT to PDF conversion.")
    yield str(file_path)


@pytest.fixture
def sample_csv(tmp_path) -> Iterator[str]:
    """Create a CSV file for testing."""
    file_path = tmp_path / f"data_{uuid.uuid4().hex}.csv"
    with open(file_path, "w") as fh:
        fh.write("Name,Age,Country\nAlice,30,USA\nBob,25,UK\nCharlie,35,Canada\n")
    yield str(file_path)


@pytest.fixture
def sample_json(tmp_path) -> Iterator[str]:
    """Create a JSON file for testing."""
    file_path = tmp_path / f"data_{uuid.uuid4().hex}.json"
    with open(file_path, "w") as fh:
        fh.write('[{"name": "Alice", "age": 30, "city": "NYC"}, {"name": "Bob", "age": 25, "city": "LA"}]')
    yield str(file_path)
